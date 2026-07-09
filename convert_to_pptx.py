"""
convert_to_pptx.py

Usage:
  1) Install dependencies: pip install python-pptx beautifulsoup4 lxml
  2) Run: python convert_to_pptx.py

This script reads `index.html` in the same folder and produces
`Portfolio_Rodrigue_Tadmbaye.pptx` with one slide per main section.
"""

import os
import sys

try:
    from bs4 import BeautifulSoup
except Exception:
    print("BeautifulSoup4 is required. Install with: pip install beautifulsoup4 lxml")
    sys.exit(1)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except Exception:
    print("python-pptx is required. Install with: pip install python-pptx")
    sys.exit(1)

HTML_FILE = 'index.html'
OUTPUT_PPTX = 'Portfolio_Rodrigue_Tadmbaye.pptx'
IMAGE_FALLBACK = 'lklklk.png'

if not os.path.exists(HTML_FILE):
    print(f"{HTML_FILE} not found in current folder: {os.getcwd()}")
    sys.exit(1)

with open(HTML_FILE, 'r', encoding='utf-8') as f:
    html = f.read()

soup = BeautifulSoup(html, 'lxml')

# Define sections in the desired order: (Display Title, element id)
sections = [
    ('Accueil', 'hero'),
    ("À propos", 'about'),
    ("Expériences", 'experience'),
    ('Projets', 'projects'),
    ('Contact', 'contact'),
]

prs = Presentation()
blank_layout = prs.slide_layouts[6]  # blank

for title_text, section_id in sections:
    section = soup.find(id=section_id)
    if section is None:
        # still add a slide with the section title
        slide = prs.slides.add_slide(blank_layout)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(32)
        p.font.bold = True
        continue

    slide = prs.slides.add_slide(blank_layout)

    # Title
    title_el = section.find(['h1','h2','h3'])
    title_str = title_el.get_text(strip=True) if title_el else title_text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_str
    p.font.size = Pt(30)
    p.font.bold = True

    # Insert hero image only on hero section (if present)
    if section_id == 'hero':
        img = section.find('img')
        img_path = None
        if img and img.get('src'):
            candidate = img.get('src')
            if os.path.exists(candidate):
                img_path = candidate
        if not img_path and os.path.exists(IMAGE_FALLBACK):
            img_path = IMAGE_FALLBACK
        if img_path:
            # place on the right
            try:
                slide.shapes.add_picture(img_path, Inches(6.4), Inches(1.6), width=Inches(2.6), height=Inches(2.6))
            except Exception as e:
                print('Warning: could not add image', e)

    # Collect textual content: paragraphs and list items
    content_top = 1.6
    left = Inches(0.6)
    width = Inches(5.8)
    height = Inches(4.5)
    content_box = slide.shapes.add_textbox(left, Inches(content_top), width, height)
    content_tf = content_box.text_frame
    content_tf.word_wrap = True

    # Prefer lists inside the section
    items = []
    # Grab direct paragraphs
    for p in section.find_all('p', recursive=True):
        txt = p.get_text(strip=True)
        if txt:
            items.append(txt)
    # Grab list items
    for li in section.find_all('li', recursive=True):
        txt = li.get_text(strip=True)
        if txt:
            items.append('• ' + txt)

    if not items:
        # fallback: take section text
        st = section.get_text(separator='\n', strip=True)
        if st:
            items = [st]

    # Add items as bullet paragraphs
    first = True
    for it in items:
        if first:
            content_tf.text = it
            first = False
            # style
            content_tf.paragraphs[0].font.size = Pt(18)
            content_tf.paragraphs[0].font.color.rgb = RGBColor(40,40,40)
        else:
            p = content_tf.add_paragraph()
            p.text = it
            p.level = 0
            p.font.size = Pt(16)

# Save presentation
prs.save(OUTPUT_PPTX)
print('Saved', OUTPUT_PPTX)
